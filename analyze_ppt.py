"""分析PPT模板结构"""
from pptx import Presentation

prs = Presentation(r"D:\pycharm\pythonProject\game\汇报.pptx")
print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Slides: {len(prs.slides)}")
print(f"Slide layouts: {len(prs.slide_layouts)}")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name}")
    for ph in layout.placeholders:
        print(f"    Placeholder idx={ph.placeholder_format.idx}: {ph.name} type={ph.placeholder_format.type}")

print()
for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i}: layout='{slide.slide_layout.name}' ===")
    for shape in slide.shapes:
        print(f"  Shape: type={shape.shape_type}, name='{shape.name}'")
        print(f"    pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})")
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                text = para.text[:100] if para.text else '(empty)'
                if j < 3:
                    print(f"    Para[{j}]: '{text}'")
        if shape.has_table:
            print(f"    TABLE: rows={shape.table.rows.__len__()}")
