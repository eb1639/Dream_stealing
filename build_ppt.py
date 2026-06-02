"""
基于模板"汇报.pptx"修改生成《盗梦空间》游戏汇报PPT
策略：清空所有shape文本，按slide索引重新填入正确内容
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE_PATH = r"D:\pycharm\pythonProject\game\汇报.pptx"
OUTPUT_PATH = r"D:\pycharm\pythonProject\game\DreamLobby_Report.pptx"

prs = Presentation(TEMPLATE_PATH)


def clear_all_text(slide):
    """清除slide中所有shape的文本"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                para.clear()


def set_shape_text(shape, lines, font_name='微软雅黑', font_size=12,
                   color=None, bold=False, alignment=None):
    """设置shape的多行文本"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.clear()

    if isinstance(lines, str):
        lines = [lines]

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            if bold:
                run.font.bold = True
        if alignment is not None:
            p.alignment = alignment


def get_text_shapes_by_pos(slide):
    """获取含文本的shapes，按top排序"""
    shapes = [(s.top, s.left, s) for s in slide.shapes if s.has_text_frame]
    shapes.sort(key=lambda x: (x[0], x[1]))
    return [s for _, _, s in shapes]


# ================================================================
# SLIDE 0: 封面
# ================================================================
slide = prs.slides[0]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
# 通常: shape[top]主标题, shape[middle]项目汇报/技术栈, shape[bottom]演示日期
if len(text_shapes) >= 1:
    set_shape_text(text_shapes[0], '《盗梦空间》', font_size=36, bold=True)
if len(text_shapes) >= 2:
    set_shape_text(text_shapes[1],
                   ['项目汇报', 'Python + Pygame'], font_size=18, bold=True)
if len(text_shapes) >= 3:
    set_shape_text(text_shapes[2],
                   ['演示日期', '2026年5月'], font_size=14)

# ================================================================
# SLIDE 1: 项目概述
# ================================================================
slide = prs.slides[1]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
# 标题 / 游戏名 / 描述 / 4个特性卡片 / 开发目标标题 / 开发目标内容
if len(text_shapes) >= 1:
    set_shape_text(text_shapes[0], ' 项目概述', font_size=28, bold=True)
if len(text_shapes) >= 2:
    set_shape_text(text_shapes[1], '《盗梦空间》 — 不对称竞技塔防游戏', font_size=18, bold=True)
if len(text_shapes) >= 3:
    set_shape_text(text_shapes[2],
        '一款基于Python + Pygame技术栈开发的像素风不对称竞技塔防游戏。'
        '玩家可选择人类或猎梦者阵营，通过建造防御建筑、管理经济资源、'
        '策略布局来对抗敌方阵营，融合了塔防、RTS和不对称竞技的核心玩法要素。'
        '游戏包含9个独立房间的矩形地图、AI自动决策系统、'
        '完整的战斗/经济/升级体系，提供深度的策略对抗体验。',
        font_size=12)
# 4个特性卡片标题 + 描述 (共8个shapes)
feature_titles = ['双阵营对抗', '策略塔防', 'AI智能系统', '像素美术风格']
feature_descs = [
    '扮演人类建造防御，或化身猎梦者突破防线，两种截然不同的游戏体验。',
    '建造门/床/炮塔/维修台，升级建筑抵御猎梦者，合理分配金币资源。',
    'AI人类自动建造升级，AI猎梦者自主选择攻击目标，实现完整的单机对战体验。',
    '纯Python程序化纹理生成，复古像素砖块/地砖/动画效果，打造沉浸式视觉体验。',
]
for i, (title, desc) in enumerate(zip(feature_titles, feature_descs)):
    idx = 4 + i * 2
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], title, font_size=14, bold=True)
    if idx + 1 < len(text_shapes):
        set_shape_text(text_shapes[idx + 1], desc, font_size=10)
# 项目开发目标
goal_title_idx = 4 + len(feature_titles) * 2
goal_content_idx = goal_title_idx + 1
if goal_title_idx < len(text_shapes):
    set_shape_text(text_shapes[goal_title_idx], '项目开发目标', font_size=16, bold=True)
if goal_content_idx < len(text_shapes):
    set_shape_text(text_shapes[goal_content_idx],
        '通过架构化构建游戏核心系统，掌握游戏开发完整流程。'
        '重点攻克复杂的状态管理逻辑、实时性能优化、'
        '以及响应不同策略的AI算法设计。'
        '最终将理论知识转化为可交互的实战作品，实现从编程思维到产品实现的完整闭环。',
        font_size=11)

# ================================================================
# SLIDE 2: 核心玩法展示
# ================================================================
slide = prs.slides[2]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
# 标题
if text_shapes:
    set_shape_text(text_shapes[0], ' 核心玩法展示', font_size=28, bold=True)

# 滑动2有很多卡片，重新填入内容
gameplay_content = [
    ('人类模式', '选择房间→上床锁定→建造防御建筑→升级门/床/炮塔→抵御猎梦者攻击→保护自己存活→击败猎梦者获胜'),
    ('猎梦者模式', '离开泉水→寻找目标→攻击门→破门入室→击杀人类→累积伤害升级→消灭所有人类获胜'),
    ('安全期(25秒)', '游戏开始后25秒安全期，人类快速进入房间上床占位，AI人类自动寻路分配房间'),
    ('经济系统', '每秒金币产出=床等级收益，建造/升级消耗金币，合理分配资源是制胜关键'),
    ('Perfect 判定', '炮塔子弹精确命中', '得分+100\n判定系数最强'),
    ('Great 判定', '猎梦者攻破门', '得分+70\n稳定推进的关键'),
    ('Good 判定', '维修台成功修复', '得分+30\n有效维持防线'),
    ('操作键位', '人类/猎梦者共用', 'WASD 移动\n鼠标 建造/升级\nESC 返回菜单\nP 暂停'),
]
for i, item in enumerate(gameplay_content):
    if i + 1 < len(text_shapes):
        if len(item) == 3:
            set_shape_text(text_shapes[i + 1],
                           [item[0], item[1], item[2]], font_size=10)
        else:
            set_shape_text(text_shapes[i + 1],
                           [item[0], item[1]], font_size=10)

# ================================================================
# SLIDE 3: 技术架构 (原"架构展示")
# ================================================================
slide = prs.slides[3]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], ' 技术架构总览', font_size=28, bold=True)

modules = [
    ('main.py', '游戏入口、主循环、状态机驱动（60FPS）'),
    ('config.py', '全局配置（地图尺寸、颜色、升级数据表）'),
    ('game_state.py', '游戏状态管理（阶段/模式/实体/建筑索引）'),
    ('map_data.py', '地图生成（9个矩形房间+走廊+墙壁+A*寻路）'),
    ('entities.py', '数据类（Human/DreamHunter/Bullet/Building）'),
    ('systems.py', '核心逻辑（经济/战斗/AI/移动/胜负判定）'),
    ('renderer.py', '渲染系统（地图/实体/建筑/特效绘制）'),
    ('effects.py', '特效系统（伤害数字/金币飘字/爆炸粒子）'),
    ('ui_menu.py', '主菜单UI（双模式选择按钮）'),
    ('ui_hud.py', 'HUD+弹出菜单（状态栏/建造/升级/暂停）'),
    ('ui_camera.py', '摄像机系统（跟随/缩放/屏幕坐标转换）'),
    ('input_router.py', '输入路由（键盘/鼠标事件分发）'),
]
for i, (name, desc) in enumerate(modules):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [name, desc], font_size=11)

# ================================================================
# SLIDE 4: 地图系统设计 (原"技术栈与项目结构")
# ================================================================
slide = prs.slides[4]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '地图系统设计', font_size=28, bold=True)

map_content = [
    ('地图规格', '33×26 矩形网格，每格32×32像素，总分辨率1056×832'),
    ('房间布局', '左组(6宽) | 走廊(2宽) | 中左(4宽) | 中心(6宽) | 走廊(2宽) | 右组(6宽)'),
    ('墙壁系统', '每个房间四边+四角均由墙壁包围，仅门位置断开'),
    ('走廊网络', '纵向走廊(cols 8-9, 23-24) + 横向走廊(rows 7-8, 16-17) 形成交通网络'),
    ('A*寻路', '多实体类型通行规则，门HP动态影响可达性'),
    ('房间模板', '9个矩形房间模板，位置/尺寸/门朝向可配置'),
    ('泉水机制', '地图边缘随机生成，猎梦者回血点'),
    ('出生点', '人类在地图上方中央出生，AI自动寻路到房间'),
]
for i, (title, desc) in enumerate(map_content):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=11)

# ================================================================
# SLIDE 5: 实体与建筑系统 (原"核心参数与常量")
# ================================================================
slide = prs.slides[5]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '实体与建筑系统', font_size=28, bold=True)

entity_data = [
    ('Human(人类)', 'id/金币/状态(游荡→上床→死亡)/房间归属/AI决策计时'),
    ('DreamHunter', '等级(Lv1-5)/累计伤害/HP/攻速/移速/泉水回血'),
    ('Bullet(子弹)', '炮塔发射/追踪猎梦者/速度360px/s'),
    ('Building(建筑)', '门(HP/等级)/床(金币产出)/炮塔(DPS/射程)/维修台(回血)'),
    ('门升级 (Lv1-5)', 'HP: 350→600→1000→1700→2800 | 费用: 0→50→130→300→650'),
    ('床升级 (Lv1-5)', '金币产出: 1→2→4→8→16/秒 | 费用: 0→35→90→230→550'),
    ('炮塔 (Lv1-4)', 'DPS: 8→14→24→40 | 射程: 5→8格 | 费用: 55→45→110→260'),
    ('维修台 (Lv1-5)', '回血: 80→140→220→340→520 HP/秒 | 费用: 35→40→100→250→550'),
    ('猎梦者升级阈值', '累计伤害: 0→350→1200→3500→8000 触发5级进化'),
]
for i, (title, desc) in enumerate(entity_data):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 6: 战斗与AI系统 (原"核心类设计 - GameData")
# ================================================================
slide = prs.slides[6]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '战斗与AI系统', font_size=28, bold=True)

combat_ai_data = [
    ('战斗流程', '猎梦者攻击相邻门→消耗攻击冷却→扣门HP→门破坏后可进入房间'),
    ('炮塔机制', '检测范围内猎梦者→冷却计时→发射追踪子弹→命中扣HP'),
    ('秒杀规则', '猎梦者触碰人类=秒杀（走廊游荡/破门后房间内均生效）'),
    ('逃跑机制', '猎梦者残血(HP<20%)→自动逃回泉水→回血→满血后重新进攻'),
    ('维修系统', '维修台持续回血: 每秒按等级恢复门HP'),
    ('AI人类决策', '每3-5秒决策: 升级床→升级门→建炮塔→建维修台→升级炮塔'),
    ('AI猎梦者决策', '选择最低等级门攻击→破门后入室追人→残血回泉水'),
    ('AI开局寻路', 'A*寻路到床→占房→上床锁定'),
    ('路径系统', 'A*算法+曼哈顿启发式+平滑像素移动'),
]
for i, (title, desc) in enumerate(combat_ai_data):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 7: 游戏状态机 (原"核心类设计 - Game")
# ================================================================
slide = prs.slides[7]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '游戏状态机设计', font_size=28, bold=True)

fsm_data = [
    ('PHASE_MENU(0)', '主菜单 → 选择人类或猎梦者模式 → 进入初始化'),
    ('PHASE_INIT(1)', '生成地图 → 创建实体 → 放置建筑 → 分配AI'),
    ('PHASE_SAFE(2)', '安全期25秒 → 人类自由移动 → 进入房间上床'),
    ('PHASE_PLAYING(3)', '正式游戏 → 战斗/AI/经济全开 → 实时对抗'),
    ('PHASE_VICTORY(4)', '胜利条件达成 → 结算画面'),
    ('PHASE_DEFEAT(5)', '失败 → 结算画面 → 返回菜单'),
    ('状态转换流', 'MENU→INIT→SAFE→PLAYING→VICTORY/DEFEAT→MENU'),
    ('双模式系统', 'MODE_HUMAN: 玩家操控人类 | MODE_HUNTER: 玩家操控猎梦者'),
    ('模式差异', '双模式共享相同地图/规则，但操作视角和策略完全不同'),
]
for i, (title, desc) in enumerate(fsm_data):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 8: 渲染与特效 (原"核心算法 - 谱面生成")
# ================================================================
slide = prs.slides[8]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '渲染与特效系统', font_size=28, bold=True)

render_data = [
    ('程序化纹理', '地板(蓝绿复古砖块)、墙壁(深棕像素砖块)、门/泉水/床'),
    ('实体渲染', '人类(像素小人多彩)、猎梦者(红色大像素体)、炮塔(金属质感)'),
    ('HUD状态栏', '顶部: 暂停按钮、金币/伤害显示、人类头像快捷栏'),
    ('弹出菜单', '建造菜单(炮塔/维修台)、升级菜单(等级/费用/属性)'),
    ('暂停菜单', '返回游戏 / 返回主菜单'),
    ('伤害飘字', '红色上飘1秒后消失，猎梦者攻击/炮塔命中触发'),
    ('金币飘字', '黄色+数字每床每秒弹出，实时反馈经济产出'),
    ('回血飘字', '绿色+数字维修时弹出，直观显示修复效果'),
    ('子弹追踪', '炮塔子弹追踪猎梦者，速度360px/s，碰撞检测销毁'),
]
for i, (title, desc) in enumerate(render_data):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 9: 摄像机与输入 (原"核心算法 - 按键判断")
# ================================================================
slide = prs.slides[9]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '摄像机与输入系统', font_size=28, bold=True)

input_data = [
    ('摄像机跟随', '平滑跟随玩家实体，支持回弹效果，始终展示最佳视角'),
    ('坐标转换', '屏幕坐标↔网格坐标，支持缩放变换，精确点击判定'),
    ('缩放适配', '动态计算地图显示区域，适配窗口大小变化'),
    ('世界边界', '摄像机限制在地图范围内，不显示外部空白区域'),
    ('键盘移动', 'WASD/方向键控制角色移动(人类4格/秒，猎梦者3.5-4格/秒)'),
    ('快捷键', 'ESC返回菜单 | P暂停/继续 | 数字键快捷操作'),
    ('鼠标交互', '点击地图建造/升级建筑 | 点击头像切换视角'),
    ('事件路由', 'InputRouter统一管理键盘/鼠标事件分发，解耦模块'),
    ('状态感知', '不同游戏阶段自动切换输入规则(菜单/游戏/暂停模式)'),
]
for i, (title, desc) in enumerate(input_data):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 10: 数据配置与平衡 (原"渲染逻辑与视觉效果")
# ================================================================
slide = prs.slides[10]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '数据配置与平衡设计', font_size=28, bold=True)

config_data = [
    ('配置项总数', '40+可调参数，涵盖游戏所有系统维度'),
    ('猎梦者5级成长', 'HP: 500→800→1200→1800→2600 | 攻击: 12→18→28→42→62'),
    ('攻速进化', '0.8→1.0→1.2→1.5→1.8 次/秒'),
    ('移速进化', '3.5→3.5→3.5→3.5→4.0 格/秒(满级提速)'),
    ('泉水回血', '60→80→100→130→170 HP/秒'),
    ('门防御梯度', 'HP: 350→600→1000→1700→2800 | 费用: 0→50→130→300→650'),
    ('床经济梯度', '金币: 1→2→4→8→16/秒 | 费用: 0→35→90→230→550'),
    ('全局参数', '人类移速4格/s | 子弹速度360px/s | AI决策3-5s | 安全期25s'),
    ('平衡理念', '前期人类时间窗口 vs 猎梦者后期爆发力，动态博弈平衡'),
]
for i, (title, desc) in enumerate(config_data):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 11: 项目亮点 (原"数据持久化")
# ================================================================
slide = prs.slides[11]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '项目亮点与技术创新', font_size=28, bold=True)

highlights = [
    ('1. 完整游戏架构', '12个模块清晰分层，状态机驱动，60FPS主循环'),
    ('2. 不对称竞技', '两大阵营截然不同的操作/策略/胜利条件'),
    ('3. 手动地图生成', '9个矩形房间精确布局，完整墙壁+2格宽走廊系统'),
    ('4. A*寻路算法', '多实体类型通行规则，门/房间状态动态更新'),
    ('5. 程序化像素美术', '纯代码生成所有纹理，无需外部图片资源'),
    ('6. 双模式AI系统', 'AI人类自动建造升级 + AI猎梦者自主攻击决策'),
    ('7. 完整战斗系统', '子弹追踪/伤害计算/升级/维修/回血'),
    ('8. 经济系统', '金币产出/消耗/升级费用/资源管理策略'),
    ('9. 流畅用户交互', '摄像机平滑跟随/缩放/弹出菜单/暂停系统'),
    ('10. 数据驱动设计', '40+可配参数，易于调整游戏平衡性'),
]
for i, (title, desc) in enumerate(highlights):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 12: 开发挑战 (原"AI在项目中的应用-编码")
# ================================================================
slide = prs.slides[12]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '开发过程与技术挑战', font_size=28, bold=True)

challenges = [
    ('地图生成算法', '从9个不规则房间到矩形布局的演进，解决墙壁与走廊空间冲突问题'),
    ('状态管理复杂度', '6个游戏阶段+2种模式+人类3状态+建筑4类型，确保状态转换正确'),
    ('AI寻路系统', '多实体类型不同通行规则，门HP动态影响可达性，房间预订防冲突'),
    ('性能优化', '60FPS渲染+AI决策+战斗计算，使用对象池/计时器/脏标记优化'),
    ('异步平衡', '人类建造速度 vs 猎梦者攻击速度的动态平衡，多轮测试调整数值'),
    ('架构挑战', '12个模块间依赖管理，避免循环导入，保持高内聚低耦合'),
    ('渲染挑战', '程序化纹理生成+像素级绘制+特效系统，无外部图片资源'),
    ('输入系统', '多阶段/多模式下的输入上下文切换，防止误操作'),
]
for i, (title, desc) in enumerate(challenges):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 13: 代码统计 (原"AI在编码中的应用-代码补全")
# ================================================================
slide = prs.slides[13]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '代码统计', font_size=28, bold=True)

stats = [
    'main.py — 410行 (游戏入口/主循环/状态机)',
    'systems.py — 583行 (经济/战斗/AI/移动/胜负)',
    'map_data.py — 330行 (地图生成+A*寻路)',
    'renderer.py — 350行 (渲染系统/程序化纹理)',
    'entities.py — 265行 (数据类/Building升级)',
    'game_state.py — 220行 (状态管理/数据查询)',
    'ui_hud.py — 200行 (HUD/建造/升级/暂停菜单)',
    'config.py — 108行 (全局配置/数据表)',
    'effects.py — 100行 (飘字/粒子特效)',
    'ui_menu.py + ui_camera.py + input_router.py — 约200行',
    '',
    '总计: 约2800行Python代码 | 12个模块 | 0外部图片资源',
]
for i, line in enumerate(stats):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], line, font_size=11)

# ================================================================
# SLIDE 14: 操作说明 (原"AI在项目中的应用-算法逻辑")
# ================================================================
slide = prs.slides[14]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '操作说明', font_size=28, bold=True)

controls = [
    ('移动操作', 'WASD / 方向键 — 控制角色移动'),
    ('建造操作', '点击房间地板 → 弹出建造菜单 → 选择建筑类型'),
    ('升级操作', '点击已有建筑 → 弹出升级菜单 → 确认升级'),
    ('视角切换', '点击顶部人类头像 → 切换到对应人类视角'),
    ('通用快捷键', 'ESC — 返回主菜单 | P — 暂停/继续'),
    ('人类模式目标', '进入房间→上床→建造防御→存活→击杀猎梦者'),
    ('猎梦者模式目标', '离开泉水→寻找目标→破门→击杀人类→升级→消灭全部'),
    ('经济策略', '优先升级床(金币产出)→升级门(防御)→建炮塔(火力)→建维修台(续航)'),
    ('生存技巧', '安全期快速占房，合理分配金币，不要过度集中防守'),
]
for i, (title, desc) in enumerate(controls):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 15: 未来展望 (原"AI在项目中的应用-文档与注释")
# ================================================================
slide = prs.slides[15]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '未来展望', font_size=28, bold=True)

future = [
    ('1. 多人联机', '引入网络对战（Socket/WebSocket），支持2-6人实时对抗'),
    ('2. 随机地图', '程序化地图生成器，增加地图多样性和可玩性'),
    ('3. 新建筑类型', '陷阱(减速/伤害)、传送门(快速移动)、隐身装置(反侦察)'),
    ('4. 音效系统', '背景音乐/战斗音效/环境音效/UI音效反馈'),
    ('5. UI美化', '更精美的像素风菜单/动画过渡/屏幕震动效果'),
    ('6. 存档系统', '游戏进度保存/排行榜/成就系统/统计数据'),
    ('7. 移动端移植', '适配触摸操作，发布Android/iOS移动版本'),
    ('8. 更多模式', '合作防守模式、无尽模式、自定义房间模式'),
]
for i, (title, desc) in enumerate(future):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 16: 游戏演示 (原"AI在项目中的应用-Bug修复")
# ================================================================
slide = prs.slides[16]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '游戏演示', font_size=28, bold=True)

demo_content = [
    '（此处插入游戏运行截图或录屏）',
    '',
    '游戏特色画面：',
    '· 主菜单双模式选择界面 — 选择人类/猎梦者阵营',
    '· 地图全貌 — 9个矩形房间+走廊+墙壁+泉水',
    '· 人类模式 — 进入房间/建造/升级/防御',
    '· 猎梦者模式 — 破门/追击/击杀/升级',
    '· 战斗特效 — 子弹追踪/伤害飘字/金币产出',
    '· 胜利/失败结算画面',
]
for i, line in enumerate(demo_content):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], line, font_size=12)

# ================================================================
# SLIDE 17: 致谢/总结 (原"项目总结")
# ================================================================
slide = prs.slides[17]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '项目总结', font_size=28, bold=True)

summary = [
    ('核心成果', '成功构建完整的游戏主循环，实现了实时动态渲染、精准按键判定、'
     'JSON数据持久化存储，确保游戏在多轮运行中的一致性与可追溯性。'),
    ('从0到1', '完成了一款可交互的游戏产品，掌握了完整的游戏开发流程。'),
    ('技术突破', '深入理解Pygame渲染机制，熟练掌握游戏主循环中的状态管理、'
     '双重缓冲渲染、碰撞检测、时间片调度等核心模式，'
     '构建了12个模块清晰分层的游戏架构。'),
    ('AI赋能', 'AI工具在代码骨架搭建、算法验证、文档注释生成方面提供了关键支持。'
     '通过人机协同，大幅度缩短了开发周期，同时保证了技术实现方案的规范性。'),
    ('项目价值', '通过架构化构建游戏核心系统，将理论知识转化为可交互的实战作品，'
     '实现从编程思维到产品实现的完整闭环。'),
]
for i, (title, desc) in enumerate(summary):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], [title, desc], font_size=10)

# ================================================================
# SLIDE 18: 致谢 (原"感谢聆听")
# ================================================================
slide = prs.slides[18]
clear_all_text(slide)
text_shapes = get_text_shapes_by_pos(slide)
if text_shapes:
    set_shape_text(text_shapes[0], '感谢聆听', font_size=36, bold=True,
                   alignment=PP_ALIGN.CENTER)

thanks_lines = [
    '',
    '《盗梦空间》— 不对称竞技塔防游戏',
    '',
    '技术栈: Python 3.12 + Pygame 2.6',
    '开发周期: 2026年5月',
    '代码量: 约2800行 | 12个模块 | 0外部图片资源',
    '',
    '谢谢观看！',
]
for i, line in enumerate(thanks_lines):
    idx = 1 + i
    if idx < len(text_shapes):
        set_shape_text(text_shapes[idx], line, font_size=14,
                       alignment=PP_ALIGN.CENTER)

# ================================================================
# 保存
# ================================================================
prs.save(OUTPUT_PATH)
print(f"PPT saved to: {OUTPUT_PATH}")
print("Done!")
